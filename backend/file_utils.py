"""
文件文本提取模块
支持: pdf, docx, pptx, txt, html, ipynb
"""
import json
import os
from pathlib import Path

# 最大提取字符数（避免超出 LLM 上下文窗口）
MAX_EXTRACT_CHARS = 8000


def extract_text(file_path: str, filename: str) -> str:
    """
    根据文件扩展名提取文本内容
    返回提取的文本字符串（最多 MAX_EXTRACT_CHARS 字符）
    """
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    try:
        if ext == "txt":
            return _read_txt(file_path)
        elif ext == "html":
            return _read_txt(file_path)
        elif ext == "md":
            return _read_txt(file_path)
        elif ext == "ipynb":
            return _read_ipynb(file_path)
        elif ext == "pdf":
            return _read_pdf(file_path)
        elif ext == "docx":
            return _read_docx(file_path)
        elif ext == "pptx":
            return _read_pptx(file_path)
        elif ext == "xlsx":
            return _read_xlsx(file_path)
        else:
            return f"[不支持的文件格式: .{ext}]"
    except Exception as e:
        return f"[文件提取失败: {str(e)}]"


def _read_txt(path: str) -> str:
    """读取纯文本文件"""
    try:
        with open(path, "r", encoding="utf-8") as f:
            content = f.read()
    except UnicodeDecodeError:
        with open(path, "r", encoding="gbk", errors="replace") as f:
            content = f.read()
    if len(content) > MAX_EXTRACT_CHARS:
        content = content[:MAX_EXTRACT_CHARS] + "\n\n[... 内容已截断，仅展示前 {} 字符 ...]".format(MAX_EXTRACT_CHARS)
    return content


def _read_ipynb(path: str) -> str:
    """提取 Jupyter Notebook 中的文本（markdown + code cells）"""
    with open(path, "r", encoding="utf-8") as f:
        nb = json.load(f)

    parts = []
    for cell in nb.get("cells", []):
        cell_type = cell.get("cell_type", "")
        source = cell.get("source", [])
        if isinstance(source, list):
            source = "".join(source)
        if not source.strip():
            continue

        if cell_type == "markdown":
            parts.append(f"[Markdown]\n{source}")
        elif cell_type == "code":
            parts.append(f"[Code]\n{source}")
        else:
            parts.append(source)

    content = "\n\n".join(parts)
    if len(content) > MAX_EXTRACT_CHARS:
        content = content[:MAX_EXTRACT_CHARS] + "\n\n[... 内容已截断 ...]"
    return content


def _read_pdf(path: str) -> str:
    """提取 PDF 文本"""
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        return "[错误: 未安装 PyPDF2 库，无法读取 PDF 文件]"

    reader = PdfReader(path)
    parts = []
    total = 0
    for page in reader.pages:
        text = page.extract_text()
        if text:
            parts.append(text)
            total += len(text)
            if total > MAX_EXTRACT_CHARS:
                break

    content = "\n\n".join(parts)
    if len(content) > MAX_EXTRACT_CHARS:
        content = content[:MAX_EXTRACT_CHARS] + "\n\n[... 内容已截断 ...]"
    return content or "[PDF 中未提取到可读文本]"


def _read_docx(path: str) -> str:
    """提取 Word 文档文本"""
    try:
        from docx import Document
    except ImportError:
        return "[错误: 未安装 python-docx 库，无法读取 docx 文件]"

    doc = Document(path)
    parts = []
    total = 0
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
            total += len(para.text)
            if total > MAX_EXTRACT_CHARS:
                break

    content = "\n".join(parts)
    if len(content) > MAX_EXTRACT_CHARS:
        content = content[:MAX_EXTRACT_CHARS] + "\n\n[... 内容已截断 ...]"
    return content or "[文档中未提取到可读文本]"


def _read_pptx(path: str) -> str:
    """提取 PowerPoint 演示文稿文本"""
    try:
        from pptx import Presentation
    except ImportError:
        return "[错误: 未安装 python-pptx 库，无法读取 pptx 文件]"

    prs = Presentation(path)
    parts = []
    total = 0
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        slide_parts.append(para.text)
                        total += len(para.text)
        if slide_parts:
            parts.append(f"--- 第 {slide_num} 页 ---\n" + "\n".join(slide_parts))
        if total > MAX_EXTRACT_CHARS:
            break

    content = "\n\n".join(parts)
    if len(content) > MAX_EXTRACT_CHARS:
        content = content[:MAX_EXTRACT_CHARS] + "\n\n[... 内容已截断 ...]"
    return content or "[演示文稿中未提取到可读文本]"


def _read_xlsx(path: str) -> str:
    """提取 Excel 表格文本（按工作表 + 制表符分隔行）"""
    try:
        from openpyxl import load_workbook
    except ImportError:
        return "[错误: 未安装 openpyxl 库，无法读取 xlsx 文件]"

    wb = load_workbook(path, data_only=True, read_only=True)
    parts = []
    total = 0
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines = [f"--- 工作表: {sheet_name} ---"]
        for row in ws.iter_rows(values_only=True):
            row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
            if row_text.strip():
                lines.append(row_text)
                total += len(row_text)
                if total > MAX_EXTRACT_CHARS:
                    break
        parts.append("\n".join(lines))
        if total > MAX_EXTRACT_CHARS:
            break
    wb.close()

    content = "\n\n".join(parts)
    if len(content) > MAX_EXTRACT_CHARS:
        content = content[:MAX_EXTRACT_CHARS] + "\n\n[... 内容已截断 ...]"
    return content or "[Excel 文件中未提取到可读数据]"
